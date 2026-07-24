"""Generate the SAP490 mentor Diagram Pack from the final canonical PNG assets.

The deck uses one overview slide per canonical diagram. Dense diagrams also receive
continuation slides containing readable crops from the exact same PNG, so the deck
and standalone assets cannot drift.
"""

from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
RENDERED = ROOT / "docs" / "diagrams" / "rendered"
PNG_DIR = RENDERED / "png"
CROP_DIR = ROOT / ".tmp" / "diagram-deck-crops"
OUTPUT = RENDERED / "SU26SAP01_GSU26SAP01_Diagram_Pack_v1_0_20260711.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TITLE_H = Inches(0.58)
FOOTER_H = Inches(0.28)
MARGIN_X = Inches(0.32)
IMAGE_TOP = Inches(0.68)
IMAGE_BOTTOM = Inches(7.16)

# Number and direction of continuation crops for diagrams that are not readable
# when reduced to a single landscape overview slide.
DETAILS = {
    "03": ("vertical", 3),
    "04": ("vertical", 3),
    "06": ("vertical", 3),
    "08": ("vertical", 2),
    "09": ("grid", 4),
    "10": ("vertical", 2),
    "11": ("vertical", 2),
    "14": ("vertical", 2),
    "17": ("vertical", 2),
    "20": ("vertical", 2),
}


def add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(MARGIN_X, Inches(0.08), SLIDE_W - 2 * MARGIN_X, TITLE_H)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.runs[0]
    run.font.name = "Aptos Display"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(25, 63, 114)


def add_footer(slide, slide_number: int, asset: str) -> None:
    box = slide.shapes.add_textbox(
        MARGIN_X, SLIDE_H - FOOTER_H, SLIDE_W - 2 * MARGIN_X, FOOTER_H
    )
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.text = f"IDTS SAP01 | Source asset: {asset}.png | Slide {slide_number}"
    paragraph.alignment = PP_ALIGN.RIGHT
    run = paragraph.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(90, 100, 115)


def add_contained_picture(slide, image_path: Path) -> None:
    with Image.open(image_path) as image:
        width_px, height_px = image.size
    area_w = SLIDE_W - 2 * MARGIN_X
    area_h = IMAGE_BOTTOM - IMAGE_TOP
    scale = min(area_w / width_px, area_h / height_px)
    width = int(width_px * scale)
    height = int(height_px * scale)
    left = int((SLIDE_W - width) / 2)
    top = int(IMAGE_TOP + (area_h - height) / 2)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def write_vertical_crops(source: Path, count: int) -> list[Path]:
    outputs: list[Path] = []
    with Image.open(source) as image:
        width, height = image.size
        overlap = int(height * 0.03)
        for index in range(count):
            start = max(0, int(index * height / count) - (overlap if index else 0))
            end = min(height, int((index + 1) * height / count) + (overlap if index + 1 < count else 0))
            target = CROP_DIR / f"{source.stem}-detail-{index + 1}.png"
            image.crop((0, start, width, end)).save(target)
            outputs.append(target)
    return outputs


def write_grid_crops(source: Path) -> list[Path]:
    outputs: list[Path] = []
    with Image.open(source) as image:
        width, height = image.size
        overlap_x = int(width * 0.025)
        overlap_y = int(height * 0.025)
        for row in range(2):
            for column in range(2):
                left = max(0, column * width // 2 - (overlap_x if column else 0))
                right = min(width, (column + 1) * width // 2 + (overlap_x if column == 0 else 0))
                top = max(0, row * height // 2 - (overlap_y if row else 0))
                bottom = min(height, (row + 1) * height // 2 + (overlap_y if row == 0 else 0))
                target = CROP_DIR / f"{source.stem}-detail-{row * 2 + column + 1}.png"
                image.crop((left, top, right, bottom)).save(target)
                outputs.append(target)
    return outputs


def add_slide(prs: Presentation, title: str, image_path: Path, asset: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(255, 255, 255)
    add_title(slide, title)
    add_contained_picture(slide, image_path)
    add_footer(slide, len(prs.slides), asset)
    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        f"Mentor review summary: {title}. "
        "Use this diagram to explain the actors/components, the direction of the flow, "
        "the CAP/backend control boundary, and the observable data or side effect. "
        f"Canonical source asset: {asset}.png; editable source and manifest are kept in docs/diagrams/rendered."
    )


def main() -> None:
    manifest = json.loads((RENDERED / "manifest.json").read_text(encoding="utf-8"))
    diagrams = manifest["diagrams"]
    if manifest.get("diagramCount") != 21 or len(diagrams) != 21:
        raise RuntimeError("Diagram manifest must contain exactly 21 canonical diagrams")

    CROP_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    while prs.slides:
        slide_id = prs.slides._sldIdLst[0]
        prs.part.drop_rel(slide_id.rId)
        del prs.slides._sldIdLst[0]

    for entry in diagrams:
        asset = entry["asset"]
        diagram_id = entry["id"]
        title = entry["title"]
        source = PNG_DIR / f"{asset}.png"
        if not source.exists() or source.stat().st_size == 0:
            raise FileNotFoundError(f"Missing rendered PNG: {source}")
        add_slide(prs, f"{diagram_id}. {title} — Overview", source, asset)

        detail = DETAILS.get(diagram_id)
        if not detail:
            continue
        mode, count = detail
        crops = write_grid_crops(source) if mode == "grid" else write_vertical_crops(source, count)
        for index, crop in enumerate(crops, start=1):
            add_slide(
                prs,
                f"{diagram_id}. {title} — Detail {index}/{len(crops)}",
                crop,
                asset,
            )

    prs.core_properties.title = "IDTS SAP01 Diagram Pack"
    prs.core_properties.subject = "SAP490 mentor review - current canonical diagrams"
    prs.core_properties.author = "IDTS SAP01 Project Team"
    prs.core_properties.comments = (
        "Generated from docs/diagrams/rendered/png; overview and detail slides use the same final assets."
    )
    prs.save(OUTPUT)
    print(f"Generated {len(prs.slides)} slides: {OUTPUT.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
